import sys

from pptx import Presentation
from pyshorteners import Shortener


def shorten_ppt_urls(file_name: str, output_file_name: str) -> None:
    """
    This function shortens all URLs in a PowerPoint presentation using the TinyURL service.

    Args:
    file_name (str): The name of the PowerPoint file to be modified.
    output_file_name (str): The name of the output PowerPoint file where the changes will be saved.

    Returns:
    None
    """

    # Open the PowerPoint file
    prs = Presentation(file_name)

    # Create a shortener object using the TinyURL service
    shortener = Shortener().tinyurl

    # Iterate through all the slides, shapes, and runs in the presentation
    for slide in prs.slides:
        for shape in slide.shapes:
            # Check if the shape has a text frame
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # Get the hyperlink address
                    address = run.hyperlink.address

                    # Check if the address is not None
                    if address is None:
                        continue

                    # Try to shorten the URL
                    try:
                        run.text = shortener.short(address)
                    except:
                        continue

    # Save the changes to a new PowerPoint file
    prs.save(file_name)


if __name__ == "__main__":
    if len(sys.argv) == 3:
        ppt_filename = sys.argv[1]
        output_ppt_filename = sys.argv[2]
    else:
        ppt_filename = input("Please provide the file name: ")
        output_ppt_filename = input("Please provide the output file name: ")

    # Use the function
    shorten_ppt_urls(ppt_filename)
