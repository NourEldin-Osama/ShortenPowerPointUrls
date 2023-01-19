from pptx import Presentation
from pyshorteners import Shortener


def shorten_ppt_urls(file_name: str) -> None:
    """
    This function shortens all URLs in a PowerPoint presentation using the TinyURL service.

    Args:
    file_name (str): The name of the PowerPoint file to be modified.

    Returns:
    None
    """

    # Open the PowerPoint file
    prs = Presentation(file_name)

    # Create a shortener object using the TinyURL service
    shortener = Shortener().tinyurl

    # Iterate through all the slides, shapes, and runs in the presentation
    for slide in prs.slides:
        for shape in slide.shapes:
            # Check if the shape has a text frame
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # Get the hyperlink address
                    address = run.hyperlink.address

                    # Check if the address is not None
                    if address is None:
                        continue

                    # Try to shorten the URL
                    try:
                        run.text = shortener.short(address)
                    except:
                        continue

    # Save the changes to the original PowerPoint file
    prs.save(file_name)


# Use the function
FILE_NAME = 'PYPPT.pptx'
shorten_ppt_urls(FILE_NAME)
